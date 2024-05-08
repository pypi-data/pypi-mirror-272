import os
from pptx import Presentation
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain.docstore.document import Document
from tqdm import tqdm

def IndexFromPPTs(folder_path: str, embeddings, vectorstore, chunk_size: int = 2000):
    """
    ### Load pdf from folder
    this function loads documents from a folder and creates a Fass index (db).\n
    - folder_path: Path to folder containing PPTs\n
    - chunk_size: Chunk size in characters (default 2000)\n
    - embeddings: Embeddings object from langchain.embeddings\n
    - vectorstore: could be FAISS or ChromaDB\n
    """

    # Assuming your PDFs are in a folder named 'test_pdfs'
    ppt_files = [f for f in os.listdir(folder_path) if f.endswith('.pptx')]

    documents = []
    # Process each PDF file and split into chunks of 24000 characters
    text_splitter = RecursiveCharacterTextSplitter(chunk_size=chunk_size, chunk_overlap=0, length_function=len, is_separator_regex=False) # FIXED

    documents = []
    # Process each PDF file
    for ppt_file in tqdm(ppt_files, "Parsing ppts"):
        prs = Presentation(ppt_file)
        total_text = ''
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    total_text += str(shape.text) # Extract text from each page
        
        texts = text_splitter.split_text(total_text)

        for text in texts:
            document = Document(page_content=text, metadata={"title": ppt_file, "type": "pptx", "tags" : ["document"]})
            documents.append(document)


    # Store Embeddings in FAISS
    for i in tqdm(range(0,1), "Creating db from documents"):
        db = vectorstore.from_documents(documents, embeddings)
    return db