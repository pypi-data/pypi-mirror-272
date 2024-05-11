import logging
import os
from multiprocessing.pool import Pool

import tqdm


def _from_pdf(file_path):
    try:
        import pypdf
    except ImportError:
        pypdf = None
        logging.warning("无法导入模块：pypdf，使用pip install pypdf来安装")
    strings = ""
    if pypdf is not None:
        reader = pypdf.PdfReader(file_path)
        for page in reader.pages:
            strings += page.extract_text()
    return strings


def _from_txt(file_path):
    strings = ""
    for encoding in ["GBK", "utf-8"]:
        try:
            with open(file_path, encoding=encoding) as f:
                strings = f.read()
            break
        except UnicodeDecodeError:
            logging.warning(f"文件：{file_path} 编码错误，尝试使用其他编码")
        except Exception as e:
            logging.error(f"文件：{file_path} 错误的将二进制文件当成文本文件处理" + str(e))
    return strings


def _from_pptx(file_path):
    try:
        import pptx
    except ImportError:
        pptx = None
        logging.warning("无法导入模块：pptx，使用pip install pypptx来安装，并确保python版本小于3.9")
    strings = ""
    if pptx is not None:
        ppt = pptx.Presentation(file_path)
        for slide in ppt.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    tf = shape.text_frame
                    for p in tf.paragraphs:
                        strings += p.text + " "
    return strings


def _from_docx(file_path):
    try:
        import docx
    except ImportError:
        docx = None
        logging.warning("无法导入模块：docx，使用pip install python-docx来安装")
    strings = ""
    if docx is not None:
        doc = docx.Document(file_path)
        for para in doc.paragraphs:
            strings += para.text
    return strings


_methods = {".pdf": _from_pdf,
            ".pptx": _from_pptx,
            ".docx": _from_docx}


class Reader:
    @staticmethod
    def read(file_path):
        name, suffix = os.path.splitext(os.path.join(file_path))
        try:
            method = _methods.get(suffix, _from_txt)
            return method(file_path)
        except UnicodeDecodeError:
            logging.warning(f"文件：{file_path} 不能正确解析。")
        except Exception as e:
            logging.error(f"读取文件时发生错误：{str(e)}")
        return ""

    @classmethod
    def read_directory(cls, directory):
        paths = []
        temp_results = []

        pool = Pool()
        for root, dirs, files in os.walk(directory):
            for file in files:
                path = os.path.normpath(os.path.join(root, file))
                paths.append(path)
                res = pool.apply_async(cls.read, args=(path,))
                temp_results.append(res)
        pool.close()
        strings = [item.get() for item in tqdm.tqdm(temp_results, desc="Reading.")]
        return strings, paths
